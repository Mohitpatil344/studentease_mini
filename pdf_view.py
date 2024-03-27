from tkinter import *
from tkinter import filedialog
from PyPDF2 import PdfReader
import os

def open_file():
    file_path = filedialog.askopenfilename(filetypes=[("Supported Files", "*.pdf;*.docx;*.pptx")])
    if file_path:
        try:
            if file_path.endswith('.pdf'):
                display_pdf(file_path)
            elif file_path.endswith('.docx'):
                display_docx(file_path)
            elif file_path.endswith('.pptx'):
                display_pptx(file_path)
            else:
                raise Exception("Unsupported file format.")
        except Exception as e:
            text.delete(1.0, END)
            text.insert(END, f"Error: {str(e)}")

def display_pdf(file_path):
    pdf_text = ""
    with open(file_path, "rb") as f:
        pdf_reader = PdfReader(f)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            pdf_text += page.extract_text() + "\n\n"  # Adding newlines between pages

    text.config(state="normal")  # Re-enable the text widget to display PDF content
    text.delete(1.0, END)
    text.insert(END, pdf_text)
    text.config(state="disabled")

def display_docx(file_path):
    import docx
    doc = docx.Document(file_path)
    text_content = ""
    for paragraph in doc.paragraphs:
        text_content += paragraph.text + "\n"
    text.config(state="normal")  # Re-enable the text widget to display docx content
    text.delete(1.0, END)
    text.insert(END, text_content)
    text.config(state="disabled")

def display_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text_content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
    text.config(state="normal")  # Re-enable the text widget to display pptx content
    text.delete(1.0, END)
    text.insert(END, text_content)
    text.config(state="disabled")

root = Tk()
root.geometry("1240x750")
root.title("File Viewer")

text = Text(root, wrap=WORD)
text.pack(expand=True, fill=BOTH)

browse_button = Button(root, text="Open File", command=open_file)
browse_button.pack()

root.mainloop()
